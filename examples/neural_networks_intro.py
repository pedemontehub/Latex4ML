"""
Presentación: Introducción a las Redes Neuronales Artificiales
Audiencia: Estudiantes de licenciatura / posgrado
Idioma: Español
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent / "templates" / "python"))
from helpers import (
    new_prs, blank, add_rect, add_text, add_figure,
    render_equation, header_bar, footer_bar,
    slide_title, slide_agenda, slide_bullets,
    slide_two_col, slide_figure, slide_thank_you,
    ML_BLUE, ML_ORANGE, ML_GREEN, ML_GRAY,
    ML_WHITE, ML_LIGHT, ML_DARK,
    SLIDE_W, SLIDE_H,
)
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
import numpy as np

OUT   = Path(__file__).parent / "neural_networks_intro.pptx"
TOTAL = 16


# ── Custom figures ────────────────────────────────────────────────────────────

def fig_motivacion():
    fig, ax = plt.subplots(figsize=(6, 3.5), facecolor="none")
    years = [2012, 2014, 2016, 2018, 2020, 2022, 2024]
    deep  = [63,   73,   82,   87,   90,   94,   97]
    svm   = [55,   62,   68,   72,   75,   None, None]
    ax.plot(years, deep, "o-",  color="#005293", lw=2.5, label="Deep NN", ms=7)
    svm_y = [v for v in svm if v]
    svm_x = years[:len(svm_y)]
    ax.plot(svm_x, svm_y, "s--", color="#DC6400", lw=2, label="SVM/RF", ms=6)
    ax.set_xlabel("Año", fontsize=12)
    ax.set_ylabel("Top-5 Accuracy (%)", fontsize=12)
    ax.set_title("ImageNet — evolución de la precisión", fontsize=13,
                 fontweight="bold")
    ax.legend(fontsize=11)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_ylim(50, 100)
    ax.grid(axis="y", alpha=0.3)
    fig.tight_layout()
    return fig


def fig_neurona():
    fig, ax = plt.subplots(figsize=(8, 4), facecolor="none")
    ax.axis("off")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)

    # Inputs
    inputs = [(0.5, 4.0), (0.5, 2.5), (0.5, 1.0)]
    labels = ["$x_1$", "$x_2$", "$x_3$"]
    weights = ["$w_1$", "$w_2$", "$w_3$"]
    for (x, y), lbl, w in zip(inputs, labels, weights):
        circle = plt.Circle((x, y), 0.35, color="#005293", zorder=3)
        ax.add_patch(circle)
        ax.text(x, y, lbl, ha="center", va="center",
                fontsize=11, color="white", fontweight="bold")
        ax.annotate("", xy=(4.5, 2.5), xytext=(x + 0.35, y),
                    arrowprops=dict(arrowstyle="->", color="#808080", lw=1.5))
        mx, my = (x + 4.5) / 2 + 0.1, (y + 2.5) / 2 + 0.15
        ax.text(mx, my, w, fontsize=10, color="#DC6400")

    # Neuron body
    neuron = plt.Circle((4.5, 2.5), 0.6, color="#005293", zorder=4)
    ax.add_patch(neuron)
    ax.text(4.5, 2.5, "$\\sigma$", ha="center", va="center",
            fontsize=14, color="white", fontweight="bold")

    # Bias
    ax.text(4.5, 1.4, "$+b$", ha="center", fontsize=12, color="#008246")

    # Output arrow
    ax.annotate("", xy=(7.5, 2.5), xytext=(5.1, 2.5),
                arrowprops=dict(arrowstyle="->", color="#005293", lw=2))
    ax.text(7.8, 2.5, "$y$", ha="left", va="center",
            fontsize=14, fontweight="bold", color="#005293")

    fig.tight_layout()
    return fig


def fig_activaciones():
    z = np.linspace(-4, 4, 300)
    fig, ax = plt.subplots(figsize=(7, 3.8), facecolor="none")
    ax.plot(z, 1 / (1 + np.exp(-z)), color="#005293", lw=2.5, label="Sigmoid $\\sigma(z)$")
    ax.plot(z, np.tanh(z),            color="#DC6400", lw=2.5, ls="--", label="Tanh")
    ax.plot(z, np.maximum(0, z),      color="#008246", lw=2.5, ls="-.", label="ReLU")
    ax.axhline(0, color="gray", lw=0.8, ls=":")
    ax.axvline(0, color="gray", lw=0.8, ls=":")
    ax.set_xlabel("$z$", fontsize=13)
    ax.set_ylabel("$a(z)$", fontsize=13)
    ax.legend(fontsize=11)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.25)
    fig.tight_layout()
    return fig


def fig_mlp():
    fig, ax = plt.subplots(figsize=(8, 4.5), facecolor="none")
    ax.axis("off")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    layer_sizes = [3, 4, 3, 2]
    layer_x = [1.0, 3.5, 6.0, 8.5]
    layer_labels = ["Entrada", "Oculta 1", "Oculta 2", "Salida"]
    positions = []
    for l, (n, x) in enumerate(zip(layer_sizes, layer_x)):
        ys = [5 * (i + 1) / (n + 1) for i in range(n)]
        positions.append(list(zip([x] * n, ys)))
        for y in ys:
            c = plt.Circle((x, y), 0.28, color="#005293", zorder=3, alpha=0.85)
            ax.add_patch(c)
        ax.text(x, 0.15, layer_labels[l], ha="center", fontsize=10,
                color="#505050")
    for l in range(len(positions) - 1):
        for (x1, y1) in positions[l]:
            for (x2, y2) in positions[l + 1]:
                ax.plot([x1 + 0.28, x2 - 0.28], [y1, y2],
                        color="#aaaaaa", lw=0.8, zorder=1)
    fig.tight_layout()
    return fig


def fig_training_curves():
    epochs = np.arange(1, 51)
    train = 1.4 * np.exp(-0.06 * epochs) + 0.18
    val   = 1.5 * np.exp(-0.05 * epochs) + 0.28 + 0.06 * np.log1p(epochs / 20)
    fig, ax = plt.subplots(figsize=(6, 3.5), facecolor="none")
    ax.plot(epochs, train, color="#005293", lw=2.5, label="Entrenamiento")
    ax.plot(epochs, val,   color="#DC6400", lw=2.5, ls="--", label="Validación")
    ax.axvline(28, color="#008246", lw=1.8, ls=":", label="Early stopping")
    ax.set_xlabel("Época", fontsize=12)
    ax.set_ylabel("Pérdida", fontsize=12)
    ax.legend(fontsize=11)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(alpha=0.25)
    fig.tight_layout()
    return fig


def fig_regularization_compare():
    x = np.linspace(-3, 3, 200)
    fig, axes = plt.subplots(1, 2, figsize=(8, 3), facecolor="none")
    # Overfitting
    np.random.seed(0)
    xd = np.linspace(-2.5, 2.5, 12)
    yd = np.sin(xd) + np.random.normal(0, 0.15, 12)
    coeffs = np.polyfit(xd, yd, 9)
    axes[0].scatter(xd, yd, color="#005293", zorder=3, s=40)
    axes[0].plot(x, np.poly1d(coeffs)(x), color="#DC6400", lw=2)
    axes[0].set_ylim(-2, 2)
    axes[0].set_title("Sobreajuste", fontsize=12)
    # Good fit
    coeffs2 = np.polyfit(xd, yd, 3)
    axes[1].scatter(xd, yd, color="#005293", zorder=3, s=40)
    axes[1].plot(x, np.poly1d(coeffs2)(x), color="#008246", lw=2)
    axes[1].set_ylim(-2, 2)
    axes[1].set_title("Ajuste correcto", fontsize=12)
    for ax in axes:
        ax.spines[["top", "right"]].set_visible(False)
        ax.grid(alpha=0.2)
    fig.tight_layout()
    return fig


# ── Slide builders ────────────────────────────────────────────────────────────

def custom_slide_backprop(prs, num):
    """Slide explaining backpropagation with equation images."""
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_LIGHT)
    header_bar(sld, "Retropropagación (Backpropagation)")

    # Left: explanation
    items = [
        "Regla de la cadena para calcular gradientes",
        "Propagación del error de salida hacia la entrada",
        "Actualización: θ ← θ − η ∇_θ L",
    ]
    for i, txt in enumerate(items):
        add_text(sld, f"• {txt}",
                 Inches(0.4), Inches(1.4) + i * Inches(0.65),
                 Inches(6.2), Inches(0.6),
                 font_size=16, color=ML_DARK)

    # Right: two equations
    eq1 = render_equation(
        r"\delta^{(l)} = \left(\mathbf{W}^{(l+1)}\right)^T \delta^{(l+1)} \odot \sigma'(z^{(l)})",
        font_size=18, color="#005293"
    )
    eq2 = render_equation(
        r"\frac{\partial \mathcal{L}}{\partial \mathbf{W}^{(l)}} = \delta^{(l)}\left(\mathbf{a}^{(l-1)}\right)^T",
        font_size=18, color="#005293"
    )
    add_figure(sld, eq1, Inches(6.6), Inches(1.5), Inches(6.5))
    add_figure(sld, eq2, Inches(6.6), Inches(3.0), Inches(6.5))

    footer_bar(sld, num, TOTAL)
    return sld


def custom_slide_loss(prs, num):
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_LIGHT)
    header_bar(sld, "Función de Pérdida")

    add_text(sld, "Clasificación — Entropía Cruzada:",
             Inches(0.4), Inches(1.35), Inches(12), Inches(0.5),
             font_size=17, bold=True, color=ML_BLUE)
    eq_ce = render_equation(
        r"\mathcal{L}(\theta) = -\frac{1}{N}\sum_{i=1}^{N}\sum_{k=1}^{K} y_{ik}\log\hat{y}_{ik}",
        font_size=22, color="#005293"
    )
    add_figure(sld, eq_ce, Inches(1.5), Inches(1.9), Inches(10))

    add_text(sld, "Regresión — Error Cuadrático Medio (MSE):",
             Inches(0.4), Inches(3.7), Inches(12), Inches(0.5),
             font_size=17, bold=True, color=ML_BLUE)
    eq_mse = render_equation(
        r"\mathcal{L}(\theta) = \frac{1}{N}\sum_{i=1}^{N}\left(y_i - \hat{y}_i\right)^2",
        font_size=22, color="#005293"
    )
    add_figure(sld, eq_mse, Inches(1.5), Inches(4.25), Inches(10))

    footer_bar(sld, num, TOTAL)
    return sld


# ── Main ─────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    # 1. Title
    slide_title(prs,
        title="Introducción a las Redes\nNeuronales Artificiales",
        subtitle="De la neurona biológica al aprendizaje profundo",
        author="Juan García",
        institution="Depto. de Ciencias de la Computación — UNAM",
        date="Abril 2026")

    # 2. Agenda
    slide_agenda(prs, [
        "Motivación",
        "La Neurona Artificial",
        "Redes Multicapa (MLP)",
        "Función de Pérdida y Backpropagation",
        "Entrenamiento y Optimización",
        "Regularización",
        "Conclusiones",
    ])

    # 3. Motivación
    slide_figure(prs, "¿Por qué Redes Neuronales?",
        fig=fig_motivacion(),
        bullets=[
            "Reconocimiento de imágenes",
            "Procesamiento de lenguaje natural",
            "Síntesis de voz y música",
            "Juegos: AlphaGo, AlphaFold",
            "Escalabilidad con datos y cómputo",
        ],
        caption="Evolución en ImageNet Top-5 accuracy",
        slide_num=3, total=TOTAL)

    # 4. La neurona artificial
    slide_figure(prs, "La Neurona Artificial",
        fig=fig_neurona(),
        caption="Modelo de neurona: suma ponderada + activación no lineal",
        slide_num=4, total=TOTAL)

    # 5. Funciones de activación
    slide_figure(prs, "Funciones de Activación",
        fig=fig_activaciones(),
        bullets=[
            "Sigmoid: salida ∈ (0, 1)",
            "Tanh: salida ∈ (−1, 1)",
            "ReLU: f(z) = max(0, z) — más popular",
            "ReLU evita el problema del gradiente evanescente",
        ],
        slide_num=5, total=TOTAL)

    # 6. MLP arquitectura
    slide_figure(prs, "Perceptrón Multicapa (MLP)",
        fig=fig_mlp(),
        bullets=[
            "Capas: entrada → ocultas → salida",
            "Forward pass: a^(l) = σ(W^(l) a^(l-1) + b^(l))",
            "Aproximador universal de funciones",
            "Profundidad > anchura (en la práctica)",
        ],
        slide_num=6, total=TOTAL)

    # 7. Función de pérdida
    custom_slide_loss(prs, 7)

    # 8. Backpropagation
    custom_slide_backprop(prs, 8)

    # 9. SGD
    slide_bullets(prs, "Descenso de Gradiente Estocástico (SGD)",
        items=[
            ("Algoritmo mini-batch SGD:", 0),
            ("1. Tomar mini-batch B ⊂ D de tamaño B", 1),
            ("2. Calcular gradiente: ĝ = (1/B) Σ ∇_θ ℓ(f_θ(x_i), y_i)", 1),
            ("3. Actualizar: θ ← θ − η ĝ", 1),
            ("Variantes populares:", 0),
            ("Momentum SGD — acumula dirección del gradiente", 1),
            ("Adam — tasa de aprendizaje adaptativa por parámetro", 1),
            ("AdamW — Adam con decaimiento de pesos correcto", 1),
        ],
        slide_num=9, total=TOTAL)

    # 10. Regularización - sobreajuste
    slide_figure(prs, "Sobreajuste vs. Ajuste Correcto",
        fig=fig_regularization_compare(),
        bullets=[
            "Sobreajuste: memorizа en lugar de generalizar",
            "Modelo demasiado complejo para los datos",
            "Alto error en test, bajo error en train",
        ],
        slide_num=10, total=TOTAL)

    # 11. Técnicas de regularización
    slide_two_col(prs, "Técnicas de Regularización",
        left_title="Regularización de pesos",
        left_items=[
            "L2 (weight decay): L += λ‖θ‖²",
            "L1: L += λ‖θ‖₁",
            "Penaliza pesos grandes",
            "Induce soluciones dispersas (L1)",
        ],
        right_title="Regularización estructural",
        right_items=[
            "Dropout: apagar neuronas al azar",
            "Batch Normalization",
            "Data Augmentation",
            "Early Stopping",
        ],
        slide_num=11, total=TOTAL)

    # 12. Curvas de entrenamiento
    slide_figure(prs, "Curvas de Entrenamiento y Early Stopping",
        fig=fig_training_curves(),
        bullets=[
            "Monitorear pérdida en validación",
            "Detener cuando val. loss deja de mejorar",
            "Guarda el modelo con mejor val. loss",
        ],
        slide_num=12, total=TOTAL)

    # 13. Comparativa resumen
    slide_two_col(prs, "Resumen: Hiperparámetros Clave",
        left_title="Arquitectura",
        left_items=[
            "Número de capas",
            "Neuronas por capa",
            "Función de activación",
            "Inicialización de pesos",
        ],
        right_title="Entrenamiento",
        right_items=[
            "Tasa de aprendizaje η",
            "Tamaño del mini-batch B",
            "Optimizador (SGD, Adam…)",
            "Parámetros de regularización",
        ],
        slide_num=13, total=TOTAL)

    # 14. Lo que viene
    slide_bullets(prs, "Más Allá del MLP",
        items=[
            ("Redes Convolucionales (CNN)", 0),
            ("Especializadas en datos con estructura espacial (imágenes)", 1),
            ("Capas convolucionales + pooling", 1),
            ("Redes Recurrentes (RNN, LSTM, GRU)", 0),
            ("Datos secuenciales: texto, series de tiempo", 1),
            ("Transformers y Mecanismo de Atención", 0),
            ("Estado del arte en NLP, visión, audio", 1),
            ("Base de GPT, BERT, Stable Diffusion", 1),
        ],
        slide_num=14, total=TOTAL)

    # 15. Conclusiones
    slide_bullets(prs, "Conclusiones",
        items=[
            ("Las RNA aproximan funciones complejas con capas no lineales", 0),
            ("El entrenamiento usa backpropagation + descenso de gradiente", 0),
            ("La regularización es esencial para generalizar", 0),
            ("Adam y variantes aceleran la convergencia en la práctica", 0),
            ("El deep learning escala bien con más datos y cómputo", 0),
        ],
        slide_num=15, total=TOTAL)

    # 16. Thank you
    slide_thank_you(prs,
        contact="jgarcia@fciencias.unam.mx",
        repo="github.com/Latex4ML/examples")

    prs.save(str(OUT))
    print(f"Guardado: {OUT}")


if __name__ == "__main__":
    build()

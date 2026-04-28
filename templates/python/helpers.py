"""
Reusable helpers for building python-pptx presentations.
Import this module from any presentation script.
"""

import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

# ── Palette ───────────────────────────────────────────────────────────────────
ML_BLUE   = RGBColor(0,   82,  147)
ML_ORANGE = RGBColor(220, 100,   0)
ML_GREEN  = RGBColor(0,   130,  70)
ML_GRAY   = RGBColor(80,   80,  80)
ML_WHITE  = RGBColor(255, 255, 255)
ML_LIGHT  = RGBColor(235, 242, 251)
ML_DARK   = RGBColor(20,   30,  50)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Presentation init ─────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slide_layouts[6]


# ── Shape helpers ─────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def add_figure(slide, fig, left, top, width, height=None, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi,
                bbox_inches="tight", transparent=True)
    buf.seek(0)
    plt.close(fig)
    if height:
        return slide.shapes.add_picture(buf, left, top, width, height)
    return slide.shapes.add_picture(buf, left, top, width)


# ── Standard slide chrome ─────────────────────────────────────────────────────

def header_bar(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0),
             SLIDE_W, Inches(1.2), fill_color=ML_BLUE)
    add_text(slide, title,
             Inches(0.4), Inches(0.12),
             Inches(12.5), Inches(0.8),
             font_size=28, bold=True, color=ML_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.88),
                 Inches(12.5), Inches(0.3),
                 font_size=13, color=RGBColor(200, 220, 255))


def footer_bar(slide, num, total):
    add_rect(slide, Inches(0), Inches(7.2),
             SLIDE_W, Inches(0.3), fill_color=ML_BLUE)
    add_text(slide, f"{num} / {total}",
             Inches(12.4), Inches(7.2),
             Inches(0.9), Inches(0.3),
             font_size=10, color=ML_WHITE, align=PP_ALIGN.RIGHT)


# ── Equation renderer ─────────────────────────────────────────────────────────

def render_equation(latex_str, font_size=22, color="white", bg=None):
    fig, ax = plt.subplots(figsize=(8, 1.2))
    ax.axis("off")
    if bg:
        fig.patch.set_facecolor(bg)
    ax.text(0.5, 0.5, f"${latex_str}$",
            ha="center", va="center",
            fontsize=font_size, color=color,
            transform=ax.transAxes)
    return fig


# ── Pre-built slides ──────────────────────────────────────────────────────────

def slide_title(prs, title, subtitle, author, institution, date):
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_DARK)
    add_rect(sld, Inches(0), Inches(5.5), SLIDE_W, Inches(2.0),
             fill_color=ML_BLUE)
    add_text(sld, title,
             Inches(0.8), Inches(1.0), Inches(11.5), Inches(2.2),
             font_size=40, bold=True, color=ML_WHITE, align=PP_ALIGN.CENTER)
    add_text(sld, subtitle,
             Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.9),
             font_size=22, color=ML_ORANGE, align=PP_ALIGN.CENTER)
    add_text(sld, f"{author}  |  {institution}",
             Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.5),
             font_size=16, color=ML_WHITE, align=PP_ALIGN.CENTER)
    add_text(sld, date,
             Inches(0.8), Inches(6.25), Inches(11.5), Inches(0.45),
             font_size=14, color=RGBColor(180, 200, 230),
             align=PP_ALIGN.CENTER)
    return sld


def slide_agenda(prs, sections):
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_LIGHT)
    header_bar(sld, "Agenda")
    for i, sec in enumerate(sections):
        y = Inches(1.4) + i * Inches(0.68)
        add_rect(sld, Inches(0.5), y + Inches(0.04),
                 Inches(0.44), Inches(0.44), fill_color=ML_BLUE)
        add_text(sld, str(i + 1),
                 Inches(0.5), y, Inches(0.44), Inches(0.52),
                 font_size=16, bold=True, color=ML_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(sld, sec,
                 Inches(1.1), y, Inches(11.5), Inches(0.52),
                 font_size=18, color=ML_DARK)
    return sld


def slide_bullets(prs, title, items, slide_num=None, total=None):
    """items: list of (text, indent_level) tuples."""
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_LIGHT)
    header_bar(sld, title)
    for i, (text, lvl) in enumerate(items):
        indent = Inches(0.6 + lvl * 0.5)
        y = Inches(1.4) + i * Inches(0.58)
        bullet = "•" if lvl == 0 else "–"
        add_text(sld, f"{bullet}  {text}",
                 indent, y, SLIDE_W - indent - Inches(0.4), Inches(0.55),
                 font_size=18 - lvl * 2,
                 color=ML_DARK if lvl == 0 else ML_GRAY)
    if slide_num and total:
        footer_bar(sld, slide_num, total)
    return sld


def slide_two_col(prs, title,
                  left_title, left_items,
                  right_title, right_items,
                  slide_num=None, total=None):
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_LIGHT)
    header_bar(sld, title)
    for col_x, col_title, col_items in [
        (Inches(0.4), left_title, left_items),
        (Inches(7.0), right_title, right_items),
    ]:
        add_rect(sld, col_x, Inches(1.4),
                 Inches(5.9), Inches(5.5), fill_color=ML_WHITE)
        add_text(sld, col_title,
                 col_x + Inches(0.2), Inches(1.55),
                 Inches(5.5), Inches(0.5),
                 font_size=16, bold=True, color=ML_BLUE)
        for i, item in enumerate(col_items):
            add_text(sld, f"• {item}",
                     col_x + Inches(0.2),
                     Inches(2.15) + i * Inches(0.56),
                     Inches(5.5), Inches(0.52),
                     font_size=15, color=ML_DARK)
    if slide_num and total:
        footer_bar(sld, slide_num, total)
    return sld


def slide_figure(prs, title, fig, bullets=None, caption="",
                 slide_num=None, total=None):
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_LIGHT)
    header_bar(sld, title)
    if bullets:
        for i, text in enumerate(bullets):
            add_text(sld, f"• {text}",
                     Inches(0.4), Inches(1.4) + i * Inches(0.6),
                     Inches(5.8), Inches(0.56),
                     font_size=16, color=ML_DARK)
        add_figure(sld, fig, Inches(6.5), Inches(1.25), Inches(6.6))
    else:
        add_figure(sld, fig, Inches(0.8), Inches(1.2), Inches(11.5))
    if caption:
        add_text(sld, caption,
                 Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
                 font_size=11, italic=True, color=ML_GRAY,
                 align=PP_ALIGN.CENTER)
    if slide_num and total:
        footer_bar(sld, slide_num, total)
    return sld


def slide_thank_you(prs, contact="", repo=""):
    sld = prs.slides.add_slide(blank(prs))
    add_rect(sld, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=ML_DARK)
    add_rect(sld, Inches(0), Inches(3.05), SLIDE_W, Inches(0.08),
             fill_color=ML_BLUE)
    add_text(sld, "¡Gracias! / Thank You!",
             Inches(0), Inches(1.2), SLIDE_W, Inches(1.5),
             font_size=48, bold=True, color=ML_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(sld, "Preguntas / Questions",
             Inches(0), Inches(3.3), SLIDE_W, Inches(0.8),
             font_size=24, color=ML_ORANGE, align=PP_ALIGN.CENTER)
    if contact:
        add_text(sld, contact,
                 Inches(0), Inches(4.5), SLIDE_W, Inches(0.5),
                 font_size=16, color=RGBColor(180, 200, 230),
                 align=PP_ALIGN.CENTER)
    if repo:
        add_text(sld, repo,
                 Inches(0), Inches(5.1), SLIDE_W, Inches(0.5),
                 font_size=16, color=RGBColor(180, 200, 230),
                 align=PP_ALIGN.CENTER)
    return sld
